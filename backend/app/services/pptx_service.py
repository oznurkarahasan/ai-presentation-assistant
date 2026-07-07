"""
PPTX text extraction, security validation, and AI-state-to-PPTX generation.
"""
from pptx import Presentation
from fastapi import UploadFile
from app.core.exceptions import FileProcessingError
from app.core.logger import logger
from app.schemas.presentation_generation import PresentationState
from app.services.file_security import clean_text, validate_item_count_and_size
from pypdf import PdfWriter, PdfReader
import io
import os
import asyncio
import subprocess  # nosec B404
import shutil
import tempfile

def validate_pptx_security(prs: Presentation, file_size: int) -> None:
    """
    Validates PPTX for security issues: slide bombs, excessive size.

    Raises:
        ValidationError: If PPTX fails security checks
    """
    validate_item_count_and_size(
        file_label="PPTX", item_label="slide", item_count=len(prs.slides), file_size=file_size
    )

async def extract_text_from_pptx(file: UploadFile, file_size: int = 0) -> tuple[list[str], str, float]:
    """
    Reads the PPTX and returns slide text with layout metadata.
    Extracts both slide text and speaker notes.
    
    Args:
        file: Uploaded PPTX file
        file_size: File size in bytes (for security validation)
        
    Returns:
        tuple[list[str], str, float]:
            - Slide texts extracted per slide
            - Orientation ('portrait' or 'landscape')
            - Aspect ratio (width / height)
    """
    try:
        # Read file content into memory
        file_content = await file.read()
        file_bytes = io.BytesIO(file_content)
        
        # Load presentation
        prs = Presentation(file_bytes)
        
        if len(prs.slides) == 0:
            raise FileProcessingError("PPTX file has no slides")
        
        # Security validation
        validate_pptx_security(prs, file_size)
        
        slides_text = []
        
        for i, slide in enumerate(prs.slides, 1):
            try:
                text_parts = []
                
                # Extract text from all shapes in slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
                
                # Extract speaker notes if available
                if slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text:
                        text_parts.append(f"Notes: {notes_frame.text}")
                
                # Combine all text from slide
                slide_text = "\n".join(text_parts) if text_parts else ""
                cleaned_text = clean_text(slide_text)
                slides_text.append(cleaned_text)
                
                logger.debug(f"Extracted slide {i}/{len(prs.slides)}")
                
            except Exception as e:
                logger.warning(f"Failed to extract slide {i}: {str(e)}")
                slides_text.append("")  # Add empty string for failed slides
        
        logger.info(f"Successfully extracted {len(slides_text)} slides from PPTX")
        
        # Get orientation
        width = prs.slide_width
        height = prs.slide_height
        orientation = "portrait" if height > width else "landscape"
        
        aspect_ratio = width / height if height > 0 else 1.777
        return slides_text, orientation, aspect_ratio

    except Exception as e:
        logger.error(f"PPTX extraction error: {str(e)}", exc_info=True)
        raise FileProcessingError(
            message="Failed to extract text from PPTX",
            details=str(e)
        )

def _strip_pdf_bookmarks(pdf_path: str) -> None:
    """Remove outline/bookmarks from PDF so the browser navpane doesn't auto-open."""
    try:
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        with open(pdf_path, "wb") as f:
            writer.write(f)
    except Exception as e:
        logger.warning(f"Failed to strip PDF bookmarks from {pdf_path}: {e}")


async def convert_to_pdf_preview(pptx_path: str) -> str | None:
    """
    Converts a PPTX file to a PDF preview using LibreOffice headless.
    Output is saved as {pptx_path}.preview.pdf next to the original file.
    Returns the preview path on success, None on failure (non-breaking).
    """
    abs_pptx = os.path.abspath(pptx_path)
    abs_outdir = os.path.dirname(abs_pptx)
    preview_path = abs_pptx + ".preview.pdf"

    def _run_conversion():
        # HOME=/tmp is often required in Docker — LibreOffice needs a writable profile dir.
        # We use tempfile.gettempdir() to be more portable while staying secure.
        env = os.environ.copy()
        env["HOME"] = tempfile.gettempdir()
        
        # Find absolute path for libreoffice to satisfy Bandit B607
        libreoffice_path = shutil.which("libreoffice") or "libreoffice"
        
        result = subprocess.run(  # nosec B603
            [
                libreoffice_path,
                "--headless",
                "--norestore",
                "--nofirststartwizard",
                "--convert-to", "pdf",
                "--outdir", abs_outdir,
                abs_pptx,
            ],
            capture_output=True,
            timeout=120,
            env=env,
        )
        return result.returncode, result.stdout.decode(errors="replace"), result.stderr.decode(errors="replace")

    try:
        loop = asyncio.get_event_loop()
        returncode, stdout, stderr = await loop.run_in_executor(None, _run_conversion)

        logger.debug(f"LibreOffice stdout: {stdout}")
        if stderr:
            logger.debug(f"LibreOffice stderr: {stderr}")

        if returncode != 0:
            logger.warning(f"LibreOffice exited with code {returncode} for {abs_pptx}. stderr: {stderr}")
            return None

        # LibreOffice creates {abs_outdir}/{basename_without_ext}.pdf — rename to preview path
        basename_no_ext = os.path.splitext(os.path.basename(abs_pptx))[0]
        libreoffice_output = os.path.join(abs_outdir, basename_no_ext + ".pdf")

        if os.path.exists(libreoffice_output):
            os.rename(libreoffice_output, preview_path)
            # Strip bookmarks/outline so the browser navpane doesn't open
            _strip_pdf_bookmarks(preview_path)
            logger.info(f"PPTX preview PDF created: {preview_path}")
            return preview_path

        logger.warning(f"LibreOffice output not found at: {libreoffice_output}. stdout: {stdout}")
    except subprocess.TimeoutExpired:
        logger.warning(f"LibreOffice conversion timed out for {abs_pptx}")
    except FileNotFoundError:
        logger.warning("LibreOffice not found. Install it to enable PPTX preview.")
    except Exception as e:
        logger.warning(f"PPTX to PDF conversion failed for {abs_pptx}: {e}")

    return None


def get_pptx_orientation(file_path: str) -> tuple[str, float]:
    """
    Quickly detects the orientation of a PPTX file.
    """
    try:
        prs = Presentation(file_path)
        width = prs.slide_width
        height = prs.slide_height
        aspect_ratio = width / height if height > 0 else 1.777
        return "portrait" if height > width else "landscape", aspect_ratio
    except Exception as e:
        logger.warning(f"Failed to detect PPTX orientation: {e}")
    return "landscape", 1.777


# --- AI-state -> PPTX export -------------------------------------------------
# Only these hosts may be fetched for PPTX image embedding. All AI-resolved
# and user-editable image URLs are expected to come from Unsplash (see
# generation_service.UNSPLASH_IMAGE_DATABASE); anything else is rejected to
# prevent SSRF against internal services / cloud metadata endpoints.
_ALLOWED_IMAGE_HOSTS = {"images.unsplash.com"}


def _is_public_ip(ip_str: str) -> bool:
    import ipaddress
    ip = ipaddress.ip_address(ip_str)
    return not (
        ip.is_private
        or ip.is_loopback
        or ip.is_link_local
        or ip.is_multicast
        or ip.is_reserved
        or ip.is_unspecified
    )


def _fetch_image_bytes_safely(url: str, timeout: int = 5):
    """Fetch an image URL for PPTX embedding, guarding against SSRF.

    Restricts fetches to an allow-list of hosts and rejects URLs that
    resolve to private/loopback/link-local addresses, so a client cannot
    point image URLs at internal services or cloud metadata endpoints.
    """
    import socket
    import urllib.request
    from urllib.parse import urlparse

    parsed = urlparse(url)
    if parsed.scheme != "https":
        raise ValueError(f"Unsupported image URL scheme: {parsed.scheme!r}")
    if parsed.hostname not in _ALLOWED_IMAGE_HOSTS:
        raise ValueError(f"Image host not allowed: {parsed.hostname!r}")

    try:
        resolved = socket.getaddrinfo(parsed.hostname, 443, proto=socket.IPPROTO_TCP)
    except socket.gaierror as exc:
        raise ValueError(f"Could not resolve image host: {parsed.hostname!r}") from exc

    for family, _, _, _, sockaddr in resolved:
        if not _is_public_ip(sockaddr[0]):
            raise ValueError(f"Image host resolves to a non-public address: {parsed.hostname!r}")

    class _NoRedirect(urllib.request.HTTPRedirectHandler):
        def redirect_request(self, req, fp, code, msg, headers, newurl):
            raise ValueError(f"Refusing to follow redirect for image URL: {newurl!r}")

    opener = urllib.request.build_opener(_NoRedirect)
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with opener.open(req, timeout=timeout) as resp:
        return resp.read()


def _hex_to_rgb(hex_color: str):
    """Convert hex color string to (r, g, b) tuple."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except (ValueError, IndexError):
        return (249, 115, 22)  # fallback: orange


def generate_pptx_from_state(state: PresentationState) -> bytes:
    """Generate a real PPTX file from PresentationState using python-pptx."""
    from concurrent.futures import ThreadPoolExecutor, as_completed
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    primary_rgb = _hex_to_rgb(state.metadata.primary_color)
    accent_rgb = _hex_to_rgb(state.metadata.accent_color)
    primary_color = RGBColor(*primary_rgb)
    accent_color = RGBColor(*accent_rgb)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # Fetch all slide images concurrently up front instead of one-by-one
    # inside the slide loop, so N images cost ~1 round-trip instead of N.
    image_urls_by_idx = {
        idx: slide_data.image.url
        for idx, slide_data in enumerate(state.slides)
        if slide_data.image and getattr(slide_data.image, "url", None)
    }
    image_bytes_by_idx: dict[int, bytes] = {}
    if image_urls_by_idx:
        with ThreadPoolExecutor(max_workers=min(8, len(image_urls_by_idx))) as pool:
            futures = {
                pool.submit(_fetch_image_bytes_safely, url): idx
                for idx, url in image_urls_by_idx.items()
            }
            for future in as_completed(futures):
                idx = futures[future]
                try:
                    image_bytes_by_idx[idx] = future.result()
                except Exception as exc:
                    logger.warning(f"Skipping image for slide {idx} during PPTX export: {exc}")

    for num_idx, slide_data in enumerate(state.slides, start=1):
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(5, 5, 7)

        # Background/side image, prefetched above (if any)
        image_bytes = image_bytes_by_idx.get(num_idx - 1)

        slide_w = Inches(13.33)
        slide_h = Inches(7.5)

        # Layout ids ("standard"/"left"/"right"/"background") must match
        # frontend/app/lib/slideLayouts.ts SLIDE_LAYOUT_IDS. Kept in sync by
        # hand since this is a separate (Python) rendering engine that can't
        # import the shared frontend module.
        if slide_data.content_type == "background" and image_bytes:
            pic_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(pic_stream, Emu(0), Emu(0), width=slide_w, height=slide_h)

            # Semi-transparent dark overlay using XML alpha
            from pptx.oxml.ns import qn
            overlay = slide.shapes.add_shape(1, Emu(0), Emu(0), slide_w, slide_h)
            overlay.line.fill.background()
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)

            sp = overlay._element
            spPr = sp.find(qn('p:spPr'))
            if spPr is None:
                spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            solid_fill = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            if solid_fill is not None:
                srgb = solid_fill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgb is not None:
                    from lxml import etree
                    alpha_el = etree.SubElement(
                        srgb,
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                    )
                    alpha_el.set('val', '45000')  # ~55% transparent

            content_left = Inches(0.8)
            content_top = Inches(1.5)
            content_width = Inches(7)
            content_height = Inches(5.5)
        elif slide_data.content_type in ("left", "right") and image_bytes:
            # Add image on one side
            img_stream = io.BytesIO(image_bytes)
            img_w = Inches(5.8)
            img_h = Inches(6.5)
            img_top = Inches(0.5)

            if slide_data.content_type == "left":
                slide.shapes.add_picture(img_stream, Inches(0.4), img_top, width=img_w, height=img_h)
                content_left = Inches(6.6)
            else:
                content_left = Inches(0.6)
                slide.shapes.add_picture(img_stream, Inches(7.1), img_top, width=img_w, height=img_h)

            content_top = Inches(1.0)
            content_width = Inches(5.8)
            content_height = Inches(6.0)
        else:
            # Standard layout: full width
            content_left = Inches(0.8)
            content_top = Inches(0.8)
            content_width = Inches(11.7)
            content_height = Inches(6.5)

        # Accent top bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), slide_w, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary_color
        bar.line.fill.background()

        # Title text box
        title_box = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_para = title_tf.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_run = title_para.add_run()
        title_run.text = slide_data.title
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(255, 255, 255)
        title_run.font.name = "Calibri"

        # Separator line below title
        sep_top = content_top + Inches(1.25)
        sep = slide.shapes.add_shape(1, content_left, sep_top, Inches(1.5), Inches(0.04))
        sep.fill.solid()
        sep.fill.fore_color.rgb = primary_color
        sep.line.fill.background()

        # Bullet items
        items_top = sep_top + Inches(0.2)
        remaining_height = content_height - Inches(1.6)
        items_box = slide.shapes.add_textbox(content_left, items_top, content_width, remaining_height)
        items_tf = items_box.text_frame
        items_tf.word_wrap = True

        for i, item in enumerate(slide_data.items):
            para = items_tf.paragraphs[0] if i == 0 else items_tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            para.space_before = Pt(6)
            # Bullet dot
            dot_run = para.add_run()
            dot_run.text = "● "
            dot_run.font.size = Pt(8)
            dot_run.font.color.rgb = accent_color
            dot_run.font.name = "Calibri"
            # Item text
            text_run = para.add_run()
            text_run.text = item
            text_run.font.size = Pt(16)
            text_run.font.color.rgb = RGBColor(200, 200, 210)
            text_run.font.name = "Calibri"

        # Slide number (bottom right)
        num_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35))
        num_tf = num_box.text_frame
        num_para = num_tf.paragraphs[0]
        num_para.alignment = PP_ALIGN.RIGHT
        num_run = num_para.add_run()
        num_run.text = str(num_idx)
        num_run.font.size = Pt(9)
        num_run.font.color.rgb = RGBColor(80, 80, 90)
        num_run.font.name = "Calibri"

        # Speaker notes
        if slide_data.speaker_note:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.speaker_note

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
