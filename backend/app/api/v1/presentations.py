from fastapi import APIRouter, Depends, UploadFile, File, Query, status, HTTPException, Request
from fastapi.responses import StreamingResponse
from openai import APIConnectionError, APIError, AuthenticationError, BadRequestError, RateLimitError
from sqlalchemy.ext.asyncio import AsyncSession
from app.api.v1 import auth
from app.core.database import get_db
from app.core.logger import logger
from app.core.exceptions import FileProcessingError, ValidationError
from app.services import pdf_service, pptx_service, embedding_service, vector_db, file_validator, generation_service
from app.core.limiter import limiter
from app.schemas.presentation_generation import PresentationGenerateRequest, PresentationGenerateResponse, PresentationState, ImageLibraryItem
from pydantic import BaseModel, Field
import asyncio
import io
import os
import shutil
import uuid
from urllib.parse import quote

router = APIRouter()

# File size limit: 50MB
MAX_FILE_SIZE = 50 * 1024 * 1024

from sqlalchemy import select
from sqlalchemy.orm import selectinload
from app.models.presentation import Presentation, PresentationSession, FileType


class PresentationTitleUpdate(BaseModel):
    title: str = Field(..., min_length=1, max_length=500)


@router.post("/generate", response_model=PresentationGenerateResponse)
@limiter.limit("5/minute")
async def generate_presentation(
    request: Request,
    payload: PresentationGenerateRequest,
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    try:
        state = await generation_service.generate_presentation_state(payload)

        slide_texts = []
        for slide in state.slides:
            parts = [slide.title]
            if slide.items:
                parts.extend(slide.items)
            if slide.speaker_note:
                parts.append(slide.speaker_note)
            slide_texts.append("\n".join([part for part in parts if part]))

        embeddings = await embedding_service.create_embeddings_batch(slide_texts) if slide_texts else []
        file_path = f"generated/ai/{uuid.uuid4().hex}.json"
        presentation = await vector_db.save_ai_presentation_with_slides(
            db=db,
            user_id=current_user.id,
            title=state.metadata.title,
            file_path=file_path,
            slide_texts=slide_texts,
            embeddings=embeddings,
            ai_content_json=state.model_dump(),
        )

        return PresentationGenerateResponse(
            presentation_id=presentation.id,
            state=state,
        )
    except AuthenticationError as exc:
        logger.error(f"OpenAI authentication failed for user {current_user.id}: {exc}")
        raise HTTPException(status_code=401, detail="OpenAI authentication failed.") from exc
    except RateLimitError as exc:
        logger.warning(f"OpenAI rate limit hit for user {current_user.id}: {exc}")
        raise HTTPException(status_code=429, detail="OpenAI rate limit exceeded. Please try again later.") from exc
    except BadRequestError as exc:
        logger.error(f"OpenAI request validation failed for user {current_user.id}: {exc}")
        raise HTTPException(status_code=400, detail="Invalid AI generation request.") from exc
    except APIConnectionError as exc:
        logger.error(f"OpenAI connection error for user {current_user.id}: {exc}")
        raise HTTPException(status_code=503, detail="OpenAI service unavailable.") from exc
    except APIError as exc:
        logger.error(f"OpenAI API error for user {current_user.id}: {exc}")
        raise HTTPException(status_code=502, detail="OpenAI service error.") from exc
    except ValueError as exc:
        logger.error(f"AI response validation failed for user {current_user.id}: {exc}")
        raise HTTPException(status_code=500, detail="AI response validation failed.") from exc
    except Exception as exc:
        logger.error(f"Presentation generation failed for user {current_user.id}: {exc}")
        raise HTTPException(status_code=500, detail="Failed to generate presentation.") from exc

@router.get("/", response_model=list)
async def list_presentations(
    db: AsyncSession = Depends(get_db),
    current_user = Depends(auth.get_current_user)
):
    stmt = select(Presentation).where(Presentation.user_id == current_user.id).order_by(Presentation.created_at.desc())
    result = await db.execute(stmt)
    presentations = result.scalars().all()
    
    return [
        {
            "id": p.id,
            "title": p.title,
            "file_name": os.path.basename(p.file_path),
            "file_path": p.file_path,
            "file_type": p.file_type.value if isinstance(p.file_type, FileType) else str(p.file_type).lower(),
            "slide_count": p.slide_count,
            "status": p.status,
            "created_at": p.created_at.isoformat() if p.created_at else None,
            "is_ai_generated": p.is_ai_generated,
        }
        for p in presentations
    ]


@router.get("/image-library", response_model=list[ImageLibraryItem])
async def get_image_library(
    current_user=Depends(auth.get_current_user),
):
    """Curated Unsplash image catalog, single source of truth shared by the
    AI generation auto-matching (generation_service.resolve_image_url) and
    the editor's manual image picker."""
    return generation_service.UNSPLASH_IMAGE_DATABASE


@router.get("/ai", response_model=list)
async def list_ai_presentations(
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    stmt = (
        select(Presentation)
        .where(Presentation.user_id == current_user.id)
        .where(Presentation.is_ai_generated.is_(True))
        .order_by(Presentation.created_at.desc())
    )
    result = await db.execute(stmt)
    presentations = result.scalars().all()

    return [
        {
            "id": p.id,
            "title": p.title,
            "file_name": os.path.basename(p.file_path),
            "file_path": p.file_path,
            "file_type": p.file_type.value if isinstance(p.file_type, FileType) else str(p.file_type).lower(),
            "slide_count": p.slide_count,
            "status": p.status,
            "created_at": p.created_at.isoformat() if p.created_at else None,
            "is_ai_generated": p.is_ai_generated,
        }
        for p in presentations
    ]


@router.get("/sessions/recent", response_model=list)
async def list_recent_sessions(
    limit: int = Query(100, ge=1, le=500),
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    stmt = (
        select(PresentationSession, Presentation)
        .join(Presentation, Presentation.id == PresentationSession.presentation_id)
        .where(PresentationSession.user_id == current_user.id)
        .order_by(PresentationSession.started_at.desc())
        .limit(limit)
    )
    result = await db.execute(stmt)
    rows = result.all()

    response = []
    for session, presentation in rows:
        duration_seconds = int(session.duration_seconds or 0)
        duration_minutes = int(duration_seconds // 60)
        response.append(
            {
                "id": session.id,
                "session_type": session.session_type.value if hasattr(session.session_type, "value") else str(session.session_type),
                "duration_seconds": duration_seconds,
                "duration_minutes": duration_minutes,
                "started_at": session.started_at.isoformat() if session.started_at else None,
                "ended_at": session.ended_at.isoformat() if session.ended_at else None,
                "presentation": {
                    "id": presentation.id,
                    "title": presentation.title,
                    "slide_count": presentation.slide_count,
                },
            }
        )

    return response


@router.delete("/sessions/{session_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_session(
    session_id: int,
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    stmt = (
        select(PresentationSession)
        .join(Presentation, Presentation.id == PresentationSession.presentation_id)
        .where(PresentationSession.id == session_id)
        .where(Presentation.user_id == current_user.id)
    )
    result = await db.execute(stmt)
    session = result.scalar_one_or_none()

    if session is None:
        raise ValidationError("Session not found")

    await db.delete(session)
    await db.commit()
    return None

@router.get("/{presentation_id}")
async def get_presentation(
    presentation_id: int,
    include_slides: bool = Query(False, description="Include slide text content in response"),
    db: AsyncSession = Depends(get_db),
    current_user = Depends(auth.get_current_user)
):
    stmt = select(Presentation).where(
        Presentation.id == presentation_id,
        Presentation.user_id == current_user.id
    )

    if include_slides:
        stmt = stmt.options(selectinload(Presentation.slides))

    result = await db.execute(stmt)
    presentation = result.scalar_one_or_none()
    
    if not presentation:
        raise ValidationError("Presentation not found")
        
    # Detect orientation and aspect ratio for frontend
    orientation = "landscape"
    aspect_ratio = 1.777
    file_type_value = presentation.file_type.value if isinstance(presentation.file_type, FileType) else str(presentation.file_type).lower()
    if file_type_value == FileType.PDF.value:
        orientation, aspect_ratio = pdf_service.get_pdf_orientation(presentation.file_path)
    elif file_type_value == FileType.PPTX.value:
        orientation, aspect_ratio = pptx_service.get_pptx_orientation(presentation.file_path)

    # Include PDF preview path for PPTX files
    pdf_preview_path = None
    if file_type_value == FileType.PPTX.value:
        preview = presentation.file_path + ".preview.pdf"
        if not os.path.exists(preview) and os.path.exists(presentation.file_path):
            # On-demand conversion for files uploaded before this feature was added
            await pptx_service.convert_to_pdf_preview(presentation.file_path)
        if os.path.exists(preview):
            pdf_preview_path = preview

    response = {
        "id": presentation.id,
        "title": presentation.title,
        "file_path": presentation.file_path,
        "file_type": file_type_value,
        "pdf_preview_path": pdf_preview_path,
        "slide_count": presentation.slide_count,
        "total_pages": presentation.slide_count,  # Added for frontend compatibility
        "status": presentation.status,
        "orientation": orientation,
        "aspect_ratio": aspect_ratio,
        "is_ai_generated": presentation.is_ai_generated,
    }

    if include_slides:
        response["slides"] = [
            {"page_number": s.page_number, "content_text": s.content_text}
            for s in presentation.slides
        ]

    return response


@router.get("/{presentation_id}/ai-state", response_model=PresentationState)
async def get_ai_presentation_state(
    presentation_id: int,
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    stmt = select(Presentation).where(
        Presentation.id == presentation_id,
        Presentation.user_id == current_user.id,
        Presentation.is_ai_generated.is_(True),
    )
    result = await db.execute(stmt)
    presentation = result.scalar_one_or_none()

    if not presentation or not presentation.ai_content_json:
        raise ValidationError("AI presentation not found")

    return PresentationState.model_validate(presentation.ai_content_json)


@router.put("/{presentation_id}/ai-state", response_model=PresentationState)
async def update_ai_presentation_state(
    presentation_id: int,
    state: PresentationState,
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    stmt = select(Presentation).where(
        Presentation.id == presentation_id,
        Presentation.user_id == current_user.id,
        Presentation.is_ai_generated.is_(True),
    )
    result = await db.execute(stmt)
    presentation = result.scalar_one_or_none()

    if not presentation:
        raise ValidationError("AI presentation not found")

    presentation.ai_content_json = state.model_dump()
    presentation.title = state.metadata.title
    presentation.slide_count = len(state.slides)
    await db.commit()
    await db.refresh(presentation)

    logger.info(f"AI presentation state updated: ID={presentation_id}, User={current_user.id}")
    return state


@router.get("/{presentation_id}/export-pptx")
async def export_ai_presentation_pptx(
    presentation_id: int,
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    stmt = select(Presentation).where(
        Presentation.id == presentation_id,
        Presentation.user_id == current_user.id,
        Presentation.is_ai_generated.is_(True),
    )
    result = await db.execute(stmt)
    presentation = result.scalar_one_or_none()

    if not presentation or not presentation.ai_content_json:
        raise ValidationError("AI presentation not found")

    state = PresentationState.model_validate(presentation.ai_content_json)

    try:
        loop = asyncio.get_event_loop()
        pptx_bytes = await loop.run_in_executor(None, _generate_pptx_from_state, state)
    except Exception as exc:
        logger.error(f"PPTX export failed for presentation {presentation_id}: {exc}")
        raise HTTPException(status_code=500, detail="Failed to generate PPTX file")

    # Content-Disposition header values must be latin-1/ASCII; the title may
    # contain non-ASCII characters (e.g. Turkish "ğ", "ş", "ı"), so provide an
    # ASCII-only fallback filename plus an RFC 5987 UTF-8 filename* for
    # browsers that support it (all modern browsers do).
    ascii_title = state.metadata.title.encode("ascii", "ignore").decode("ascii")
    safe_title = "".join(c if c.isalnum() or c in " -_" else "_" for c in ascii_title).strip("_ ") or "presentation"
    filename_ascii = f"{safe_title[:60]}.pptx"
    filename_utf8 = quote(f"{state.metadata.title[:60]}.pptx")

    return StreamingResponse(
        io.BytesIO(pptx_bytes),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={
            "Content-Disposition": (
                f'attachment; filename="{filename_ascii}"; filename*=UTF-8\'\'{filename_utf8}'
            )
        },
    )


# Only these hosts may be fetched for PPTX image embedding. All AI-resolved
# and user-editable image URLs are expected to come from Unsplash (see
# generation_service.UNSPLASH_IMAGE_DATABASE); anything else is rejected to
# prevent SSRF against internal services / cloud metadata endpoints.
_ALLOWED_IMAGE_HOSTS = {"images.unsplash.com"}


def _is_public_ip(ip_str: str) -> bool:
    import ipaddress
    ip = ipaddress.ip_address(ip_str)
    return not (
        ip.is_private
        or ip.is_loopback
        or ip.is_link_local
        or ip.is_multicast
        or ip.is_reserved
        or ip.is_unspecified
    )


def _fetch_image_bytes_safely(url: str, timeout: int = 5):
    """Fetch an image URL for PPTX embedding, guarding against SSRF.

    Restricts fetches to an allow-list of hosts and rejects URLs that
    resolve to private/loopback/link-local addresses, so a client cannot
    point image URLs at internal services or cloud metadata endpoints.
    """
    import socket
    import urllib.request
    from urllib.parse import urlparse

    parsed = urlparse(url)
    if parsed.scheme != "https":
        raise ValueError(f"Unsupported image URL scheme: {parsed.scheme!r}")
    if parsed.hostname not in _ALLOWED_IMAGE_HOSTS:
        raise ValueError(f"Image host not allowed: {parsed.hostname!r}")

    try:
        resolved = socket.getaddrinfo(parsed.hostname, 443, proto=socket.IPPROTO_TCP)
    except socket.gaierror as exc:
        raise ValueError(f"Could not resolve image host: {parsed.hostname!r}") from exc

    for family, _, _, _, sockaddr in resolved:
        if not _is_public_ip(sockaddr[0]):
            raise ValueError(f"Image host resolves to a non-public address: {parsed.hostname!r}")

    class _NoRedirect(urllib.request.HTTPRedirectHandler):
        def redirect_request(self, req, fp, code, msg, headers, newurl):
            raise ValueError(f"Refusing to follow redirect for image URL: {newurl!r}")

    opener = urllib.request.build_opener(_NoRedirect)
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with opener.open(req, timeout=timeout) as resp:
        return resp.read()


def _hex_to_rgb(hex_color: str):
    """Convert hex color string to (r, g, b) tuple."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 3:
        hex_color = "".join(c * 2 for c in hex_color)
    try:
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
    except (ValueError, IndexError):
        return (249, 115, 22)  # fallback: orange


def _generate_pptx_from_state(state: PresentationState) -> bytes:
    """Generate a real PPTX file from PresentationState using python-pptx."""
    from concurrent.futures import ThreadPoolExecutor, as_completed
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = PptxPresentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    primary_rgb = _hex_to_rgb(state.metadata.primary_color)
    accent_rgb = _hex_to_rgb(state.metadata.accent_color)
    primary_color = RGBColor(*primary_rgb)
    accent_color = RGBColor(*accent_rgb)

    blank_layout = prs.slide_layouts[6]  # blank layout

    # Fetch all slide images concurrently up front instead of one-by-one
    # inside the slide loop, so N images cost ~1 round-trip instead of N.
    image_urls_by_idx = {
        idx: slide_data.image.url
        for idx, slide_data in enumerate(state.slides)
        if slide_data.image and getattr(slide_data.image, "url", None)
    }
    image_bytes_by_idx: dict[int, bytes] = {}
    if image_urls_by_idx:
        with ThreadPoolExecutor(max_workers=min(8, len(image_urls_by_idx))) as pool:
            futures = {
                pool.submit(_fetch_image_bytes_safely, url): idx
                for idx, url in image_urls_by_idx.items()
            }
            for future in as_completed(futures):
                idx = futures[future]
                try:
                    image_bytes_by_idx[idx] = future.result()
                except Exception as exc:
                    logger.warning(f"Skipping image for slide {idx} during PPTX export: {exc}")

    for num_idx, slide_data in enumerate(state.slides, start=1):
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(5, 5, 7)

        # Background/side image, prefetched above (if any)
        image_bytes = image_bytes_by_idx.get(num_idx - 1)

        slide_w = Inches(13.33)
        slide_h = Inches(7.5)

        # Layout ids ("standard"/"left"/"right"/"background") must match
        # frontend/app/lib/slideLayouts.ts SLIDE_LAYOUT_IDS. Kept in sync by
        # hand since this is a separate (Python) rendering engine that can't
        # import the shared frontend module.
        if slide_data.content_type == "background" and image_bytes:
            pic_stream = io.BytesIO(image_bytes)
            slide.shapes.add_picture(pic_stream, Emu(0), Emu(0), width=slide_w, height=slide_h)

            # Semi-transparent dark overlay using XML alpha
            from pptx.oxml.ns import qn
            overlay = slide.shapes.add_shape(1, Emu(0), Emu(0), slide_w, slide_h)
            overlay.line.fill.background()
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)

            sp = overlay._element
            spPr = sp.find(qn('p:spPr'))
            if spPr is None:
                spPr = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            solid_fill = sp.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
            if solid_fill is not None:
                srgb = solid_fill.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
                if srgb is not None:
                    from lxml import etree
                    alpha_el = etree.SubElement(
                        srgb,
                        '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha'
                    )
                    alpha_el.set('val', '45000')  # ~55% transparent

            content_left = Inches(0.8)
            content_top = Inches(1.5)
            content_width = Inches(7)
            content_height = Inches(5.5)
        elif slide_data.content_type in ("left", "right") and image_bytes:
            # Add image on one side
            img_stream = io.BytesIO(image_bytes)
            img_w = Inches(5.8)
            img_h = Inches(6.5)
            img_top = Inches(0.5)

            if slide_data.content_type == "left":
                slide.shapes.add_picture(img_stream, Inches(0.4), img_top, width=img_w, height=img_h)
                content_left = Inches(6.6)
            else:
                content_left = Inches(0.6)
                slide.shapes.add_picture(img_stream, Inches(7.1), img_top, width=img_w, height=img_h)

            content_top = Inches(1.0)
            content_width = Inches(5.8)
            content_height = Inches(6.0)
        else:
            # Standard layout: full width
            content_left = Inches(0.8)
            content_top = Inches(0.8)
            content_width = Inches(11.7)
            content_height = Inches(6.5)

        # Accent top bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), slide_w, Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = primary_color
        bar.line.fill.background()

        # Title text box
        title_box = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(1.2))
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_para = title_tf.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_run = title_para.add_run()
        title_run.text = slide_data.title
        title_run.font.size = Pt(32)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(255, 255, 255)
        title_run.font.name = "Calibri"

        # Separator line below title
        sep_top = content_top + Inches(1.25)
        sep = slide.shapes.add_shape(1, content_left, sep_top, Inches(1.5), Inches(0.04))
        sep.fill.solid()
        sep.fill.fore_color.rgb = primary_color
        sep.line.fill.background()

        # Bullet items
        items_top = sep_top + Inches(0.2)
        remaining_height = content_height - Inches(1.6)
        items_box = slide.shapes.add_textbox(content_left, items_top, content_width, remaining_height)
        items_tf = items_box.text_frame
        items_tf.word_wrap = True

        for i, item in enumerate(slide_data.items):
            para = items_tf.paragraphs[0] if i == 0 else items_tf.add_paragraph()
            para.alignment = PP_ALIGN.LEFT
            para.space_before = Pt(6)
            # Bullet dot
            dot_run = para.add_run()
            dot_run.text = "● "
            dot_run.font.size = Pt(8)
            dot_run.font.color.rgb = accent_color
            dot_run.font.name = "Calibri"
            # Item text
            text_run = para.add_run()
            text_run.text = item
            text_run.font.size = Pt(16)
            text_run.font.color.rgb = RGBColor(200, 200, 210)
            text_run.font.name = "Calibri"

        # Slide number (bottom right)
        num_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35))
        num_tf = num_box.text_frame
        num_para = num_tf.paragraphs[0]
        num_para.alignment = PP_ALIGN.RIGHT
        num_run = num_para.add_run()
        num_run.text = str(num_idx)
        num_run.font.size = Pt(9)
        num_run.font.color.rgb = RGBColor(80, 80, 90)
        num_run.font.name = "Calibri"

        # Speaker notes
        if slide_data.speaker_note:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data.speaker_note

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


@router.post("/upload", status_code=status.HTTP_201_CREATED)
async def upload_presentation(
    file: UploadFile = File(...),
    db: AsyncSession = Depends(get_db),
    current_user = Depends(auth.get_current_user)
):
    logger.info(f"Upload request from user {current_user.id}: {file.filename}")

    # Validate file extension
    if not (file.filename.endswith(".pdf") or file.filename.endswith(".pptx")):
        logger.warning(f"Invalid file type attempted: {file.filename}")
        raise ValidationError("Only PDF and PPTX files are accepted.")

    # Validate file size first (more efficient to fail early)
    file.file.seek(0, 2)  # Move to end
    file_size = file.file.tell()
    file.file.seek(0)  # Reset to beginning
    
    if file_size > MAX_FILE_SIZE:
        logger.warning(f"File too large: {file_size} bytes from user {current_user.id}")
        raise ValidationError(f"File size exceeds limit. Maximum allowed: {MAX_FILE_SIZE // (1024*1024)}MB")
    
    if file_size == 0:
        logger.warning(f"Empty file uploaded: {file.filename}")
        raise ValidationError("File is empty.")
    
    # Read first 512 bytes for magic byte validation
    file_header = await file.read(512)
    file.file.seek(0)
    
    # Validate file type using magic bytes (not just extension)
    file_validator.validate_file_type(file_header, file.filename)

    upload_dir = "uploaded_files"
    os.makedirs(upload_dir, exist_ok=True)
    
    # Generate unique filename to prevent overwrite
    import re
    unique_id = uuid.uuid4().hex
    # Sanitize the filename: remove path traversal characters and keep only safe ones
    filename_from_user = os.path.basename(file.filename)
    filename_sanitized = re.sub(r'[^a-zA-Z0-9._-]', '_', filename_from_user)
    safe_filename = f"{current_user.id}_{unique_id}_{filename_sanitized}"
    file_path = os.path.join(upload_dir, safe_filename)
    
    try:
        # Save file
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        logger.info(f"File saved: {file_path}")
        
        # Calculate file hash for analytics (optional)
        file_hash = file_validator.calculate_file_hash(file_path)
        
        file.file.seek(0)

        # Extract text based on file type (with security validation)
        if file.filename.endswith(".pdf"):
            slide_texts, _orientation, _aspect_ratio = await pdf_service.extract_text_from_pdf(file, file_size)
            logger.info(f"Extracted {len(slide_texts)} slides from PDF")
        elif file.filename.endswith(".pptx"):
            slide_texts, _orientation, _aspect_ratio = await pptx_service.extract_text_from_pptx(file, file_size)
            logger.info(f"Extracted {len(slide_texts)} slides from PPTX")
            # Convert to PDF preview for browser display (non-breaking — failure is logged only)
            await pptx_service.convert_to_pdf_preview(file_path)

        else:
            raise ValidationError("Unsupported file type")

        # Generate embeddings in parallel (10x faster!)
        logger.info(f"Generating embeddings for {len(slide_texts)} slides...")
        embeddings = await embedding_service.create_embeddings_batch(slide_texts)

        new_presentation = await vector_db.save_presentation_with_slides(
            db=db,
            user_id=current_user.id,
            title=file.filename,
            file_path=file_path,
            slide_texts=slide_texts,
            embeddings=embeddings,
            file_hash=file_hash
        )
        
        logger.info(f"Presentation uploaded successfully: ID={new_presentation.id}, User={current_user.id}")

        preview_path = file_path + ".preview.pdf"
        pdf_preview_path = preview_path if os.path.exists(preview_path) else None

        return {
            "id": new_presentation.id,
            "title": new_presentation.title,
            "pages": len(slide_texts),
            "status": "success",
            "pdf_preview_path": pdf_preview_path,
        }

    except Exception as e:
        # Clean up uploaded file on error
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"Cleaned up file after error: {file_path}")
            except Exception as cleanup_error:
                logger.warning(
                    f"Failed to clean up file after error: {file_path}. Cleanup error: {cleanup_error}",
                    exc_info=True,
                )
        
        logger.error(f"Upload failed for user {current_user.id}: {str(e)}", exc_info=True)
        raise FileProcessingError(
            message="Failed to process presentation",
            details=str(e)
        )


@router.patch("/{presentation_id}")
async def update_presentation_title(
    presentation_id: int,
    payload: PresentationTitleUpdate,
    db: AsyncSession = Depends(get_db),
    current_user=Depends(auth.get_current_user),
):
    stmt = select(Presentation).where(
        Presentation.id == presentation_id,
        Presentation.user_id == current_user.id,
    )
    result = await db.execute(stmt)
    presentation = result.scalar_one_or_none()

    if not presentation:
        raise ValidationError("Presentation not found")

    normalized_title = payload.title.strip()
    if not normalized_title:
        raise ValidationError("Presentation title cannot be empty")

    presentation.title = normalized_title
    await db.commit()
    await db.refresh(presentation)

    logger.info(f"Presentation title updated: ID={presentation.id}, User={current_user.id}")

    return {
        "id": presentation.id,
        "title": presentation.title,
        "file_name": os.path.basename(presentation.file_path),
        "file_path": presentation.file_path,
        "file_type": presentation.file_type.value if isinstance(presentation.file_type, FileType) else str(presentation.file_type).lower(),
        "slide_count": presentation.slide_count,
        "status": presentation.status,
        "created_at": presentation.created_at.isoformat() if presentation.created_at else None,
        "is_ai_generated": presentation.is_ai_generated,
    }

@router.delete("/{presentation_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_presentation(
    presentation_id: int,
    db: AsyncSession = Depends(get_db),
    current_user = Depends(auth.get_current_user)
):
    """Delete a presentation"""
    
    stmt = select(Presentation).where(
        Presentation.id == presentation_id,
        Presentation.user_id == current_user.id
    )
    result = await db.execute(stmt)
    presentation = result.scalar_one_or_none()
    
    if not presentation:
        raise ValidationError("Presentation not found")
    
    # Delete file from disk
    if os.path.exists(presentation.file_path):
        try:
            os.remove(presentation.file_path)
            logger.info(f"Deleted file: {presentation.file_path}")
        except Exception as e:
            logger.warning(f"Failed to delete file: {presentation.file_path}. Error: {e}")

    # Also delete PDF preview if it exists (PPTX uploads)
    preview_path = presentation.file_path + ".preview.pdf"
    if os.path.exists(preview_path):
        try:
            os.remove(preview_path)
            logger.info(f"Deleted PDF preview: {preview_path}")
        except Exception as e:
            logger.warning(f"Failed to delete PDF preview: {preview_path}. Error: {e}")
    
    # Delete from database (cascade will handle related records)
    await db.delete(presentation)
    await db.commit()
    
    logger.info(f"Presentation deleted: ID={presentation_id}, User={current_user.id}")
    return None